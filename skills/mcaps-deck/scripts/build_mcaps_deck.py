#!/usr/bin/env python3
"""MCAPS deck builder.

Usage:
  python3 build_mcaps_deck.py outline.json output.pptx
  python3 build_mcaps_deck.py   # writes a sample outline JSON to stdout
"""
from __future__ import annotations
import json, math, sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.dml import MSO_THEME_COLOR

W, H = 13.333, 7.5
BG_DARK = RGBColor(0x0A,0x1A,0x2F)
BG_LIGHT = RGBColor(0xF5,0xF2,0xEF)
CARD_DARK = RGBColor(0x14,0x28,0x42)
HEADER = RGBColor(0x8D,0xC8,0xE8)
PINK = RGBColor(0xE9,0x4B,0x89)
MAGENTA = RGBColor(0xA0,0x2B,0x93)
TEAL = RGBColor(0x49,0xC5,0xB1)
AZURE = RGBColor(0x00,0x78,0xD4)
WHITE = RGBColor(0xFF,0xFF,0xFF)
MUTED = RGBColor(0xB9,0xC7,0xD8)
INK = RGBColor(0x16,0x25,0x36)


def I(x): return Inches(x)
def P(x): return Pt(x)

def set_fill(shape, color):
    shape.fill.solid(); shape.fill.fore_color.rgb = color

def no_line(shape): shape.line.fill.background()

def add_bg(slide, color):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, I(W), I(H))
    set_fill(s, color); no_line(s); return s

def add_text(slide, text, x, y, w, h, size=24, color=WHITE, bold=False, font='Arial', align=None, valign=None):
    box = slide.shapes.add_textbox(I(x), I(y), I(w), I(h))
    tf = box.text_frame; tf.clear(); tf.margin_left = I(0.04); tf.margin_right = I(0.04); tf.word_wrap = True
    if valign: tf.vertical_anchor = valign
    p = tf.paragraphs[0]; p.text = str(text or '')
    if align: p.alignment = align
    r = p.runs[0] if p.runs else p.add_run()
    r.font.name = font; r.font.size = P(size); r.font.bold = bold; r.font.color.rgb = color
    return box

def add_bullets(slide, bullets, x, y, w, h, size=18, color=WHITE):
    box = slide.shapes.add_textbox(I(x), I(y), I(w), I(h))
    tf = box.text_frame; tf.clear(); tf.word_wrap = True
    for i,b in enumerate(bullets or []):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = str(b); p.level = 0; p.font.size = P(size); p.font.color.rgb = color; p.space_after = P(8)
    return box

def add_footer(slide, dark=True):
    c = MUTED if dark else RGBColor(0x65,0x68,0x70)
    add_text(slide, 'Microsoft MCAPS Israel', .55, 7.05, 3.2, .25, 8.5, c)
    add_text(slide, 'Confidential', 11.45, 7.05, 1.25, .25, 8.5, c, align=PP_ALIGN.RIGHT)

def ribbon(slide, style='blue'):
    colors = [TEAL, AZURE, MAGENTA, PINK] if style != 'purple' else [MAGENTA, PINK, TEAL, AZURE]
    for i,c in enumerate(colors):
        shape = slide.shapes.add_shape(MSO_SHAPE.PARALLELOGRAM, I(8.0+i*.55), I(.55+i*.18), I(3.6), I(.75))
        set_fill(shape,c); no_line(shape)
        shape.rotation = -10

def pill(slide, text, x, y, w=2.4):
    # simulate gradient with three overlapped rounded rects
    seg=w/3
    for i,c in enumerate([TEAL, MAGENTA, PINK]):
        sh=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, I(x+i*seg), I(y), I(seg+.08), I(.38))
        set_fill(sh,c); no_line(sh)
    add_text(slide,text,x+.08,y+.08,w-.16,.18,8.5,WHITE,True,align=PP_ALIGN.CENTER)

def card(slide,x,y,w,h,fill=CARD_DARK,line=None,radius=True):
    sh=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,I(x),I(y),I(w),I(h))
    set_fill(sh,fill)
    if line: sh.line.color.rgb=line; sh.line.width=P(1)
    else: no_line(sh)
    return sh

def title_slide(prs, s):
    slide=prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide,BG_DARK); ribbon(slide,s.get('hero_style','blue'))
    add_text(slide,'MICROSOFT AZURE',.65,.55,2.4,.28,10,MUTED,True)
    add_text(slide,s.get('title','Deck title'),.65,1.45,7.4,1.5,35,WHITE,True,'Arial Black')
    add_text(slide,s.get('subtitle',''),.72,3.05,6.7,.7,17,MUTED)
    presenters=s.get('presenters','')
    if presenters: add_text(slide,presenters,.72,5.85,5.2,.35,12,WHITE,True)
    add_footer(slide,True)

def section_divider(prs,s):
    slide=prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide,BG_DARK); ribbon(slide,'purple')
    if s.get('badge'): pill(slide,s['badge'],5.35,1.35,2.6)
    add_text(slide,s.get('title','Section'),2.0,2.5,9.3,.8,38,WHITE,True,'Arial Black',PP_ALIGN.CENTER)
    add_text(slide,s.get('subtitle',''),2.7,3.5,8.0,.45,15,MUTED,False,align=PP_ALIGN.CENTER)
    add_footer(slide,True)

def section_word(prs,s):
    slide=prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide,BG_LIGHT)
    orb=slide.shapes.add_shape(MSO_SHAPE.OVAL,I(9.6),I(1.0),I(2.5),I(2.5)); set_fill(orb,TEAL); no_line(orb)
    orb2=slide.shapes.add_shape(MSO_SHAPE.OVAL,I(10.2),I(1.4),I(2.2),I(2.2)); set_fill(orb2,PINK); no_line(orb2)
    add_text(slide,s.get('title','Build'),.8,2.35,8.5,1.1,58,MAGENTA,True,'Arial Black')
    add_text(slide,s.get('subtitle',''),.9,3.65,7.2,.65,18,INK)
    add_footer(slide,False)

def agenda(prs,s):
    slide=prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide,BG_DARK)
    add_text(slide,'Agenda',.65,.75,4,.6,30,WHITE,True,'Arial Black')
    add_text(slide,s.get('subtitle',''),.7,1.35,6,.4,14,MUTED)
    items=s.get('items',[])[:6]
    for i,it in enumerate(items):
        y=2.05+i*.72; card(slide,1.0,y,11.2,.52,RGBColor(0x0F,0x22,0x3B),RGBColor(0x22,0x48,0x70))
        add_text(slide,f'{i+1:02}',1.18,y+.13,.45,.18,11,HEADER,True)
        add_text(slide,it,1.75,y+.1,9.8,.25,15,WHITE,True)
    add_footer(slide,True)

def feature_pill(prs,s):
    slide=prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide,BG_DARK)
    pill(slide,s.get('badge','Public preview'),.75,.75,2.7)
    add_text(slide,s.get('title','Feature headline'),.75,1.35,8.8,.8,30,WHITE,True,'Arial Black')
    add_text(slide,s.get('subtitle',''),.8,2.15,9.2,.45,14,MUTED)
    bullets=s.get('bullets',[])[:3]
    for i,b in enumerate(bullets):
        x=.75+i*4.05; card(slide,x,3.2,3.55,2.1,CARD_DARK,RGBColor(0x23,0x4A,0x6D))
        add_text(slide,str(b),x+.28,3.55,3.0,1.1,18,WHITE,True)
    add_footer(slide,True)

def three_cards(prs,s):
    slide=prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide,BG_DARK)
    add_text(slide,s.get('title','Three cards'),.65,.65,8.2,.55,27,WHITE,True,'Arial Black')
    add_text(slide,s.get('subtitle',''),.7,1.2,8.5,.35,13,MUTED)
    for i,c in enumerate((s.get('cards') or [])[:3]):
        x=.72+i*4.05; card(slide,x,2.05,3.55,3.9,CARD_DARK,RGBColor(0x25,0x4A,0x6E))
        strip=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,I(x),I(2.05),I(3.55),I(.55)); set_fill(strip,HEADER); no_line(strip)
        add_text(slide,c.get('title',''),x+.22,2.18,3.05,.22,12,INK,True)
        add_text(slide,c.get('body',''),x+.25,2.9,3.02,1.75,14,WHITE)
    add_footer(slide,True)

def pipeline(prs,s):
    slide=prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide,BG_LIGHT)
    add_text(slide,s.get('title','Pipeline'),.65,.65,8.5,.5,27,INK,True,'Arial Black')
    add_text(slide,s.get('subtitle',''),.7,1.2,8.5,.35,13,RGBColor(0x60,0x5E,0x5C))
    steps=s.get('steps',[])[:5]; n=max(1,len(steps)); gap=11.7/n
    for i,st in enumerate(steps):
        x=.75+i*gap; y=3.0
        circ=slide.shapes.add_shape(MSO_SHAPE.OVAL,I(x+.35),I(y-.65),I(.72),I(.72)); set_fill(circ,AZURE if i%2==0 else TEAL); no_line(circ)
        add_text(slide,st.get('icon',str(i+1)),x+.43,y-.43,.52,.18,13,WHITE,True,align=PP_ALIGN.CENTER)
        add_text(slide,st.get('title',''),x,y+.2,gap-.35,.35,14,INK,True,align=PP_ALIGN.CENTER)
        add_text(slide,st.get('body',''),x,y+.72,gap-.35,.8,10,RGBColor(0x4A,0x4A,0x4A),align=PP_ALIGN.CENTER)
        if i<n-1:
            arr=slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,I(x+gap-0.45),I(y-.38),I(.55),I(.22)); set_fill(arr,MAGENTA); no_line(arr)
    add_footer(slide,False)

def kpi(prs,s):
    slide=prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide,BG_DARK)
    add_text(slide,s.get('title','KPI snapshot'),.65,.65,8.8,.55,27,WHITE,True,'Arial Black')
    stats=s.get('stats',[])[:4]
    for i,st in enumerate(stats):
        x=.75+i*3.1; card(slide,x,2.1,2.7,3.2,CARD_DARK,RGBColor(0x25,0x4A,0x6E))
        bar=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,I(x),I(2.1),I(.08),I(3.2)); set_fill(bar,PINK); no_line(bar)
        val=str(st.get('value','')); size=32 if len(val)<8 else 23
        add_text(slide,val,x+.25,2.55,2.15,.7,size,WHITE,True,'Arial Black')
        add_text(slide,st.get('label',''),x+.27,3.35,2.15,.45,13,HEADER,True)
        add_text(slide,st.get('note',''),x+.27,4.05,2.05,.7,10,MUTED)
    add_footer(slide,True)

def quote(prs,s):
    slide=prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide,BG_DARK)
    add_text(slide,'“',.85,1.0,1.0,.9,80,PINK,True,'Arial Black')
    add_text(slide,s.get('text','Quote'),1.7,1.65,9.9,1.9,28,WHITE,True,'Arial')
    add_text(slide,s.get('attribution',''),1.75,4.15,6,.35,14,MUTED)
    add_footer(slide,True)

def two_col(prs,s):
    slide=prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide,BG_DARK)
    add_text(slide,s.get('title','Two columns'),.65,.65,8.5,.55,27,WHITE,True,'Arial Black')
    add_text(slide,s.get('subtitle',''),.7,1.2,8.5,.35,13,MUTED)
    for i,col in enumerate((s.get('columns') or [])[:2]):
        x=.8+i*6.0; card(slide,x,2.0,5.45,3.9,CARD_DARK,RGBColor(0x25,0x4A,0x6E))
        add_text(slide,col.get('title',''),x+.35,2.35,4.8,.35,18,HEADER,True)
        add_bullets(slide,col.get('bullets',[]),x+.5,3.0,4.65,2.35,14,WHITE)
    add_footer(slide,True)

def bullets(prs,s):
    slide=prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide,BG_DARK)
    add_text(slide,s.get('title','Bullets'),.65,.65,8.5,.55,27,WHITE,True,'Arial Black')
    add_text(slide,s.get('subtitle',''),.7,1.2,8.5,.35,13,MUTED)
    add_bullets(slide,s.get('bullets',[]),1.0,2.0,10.8,3.8,19,WHITE)
    add_footer(slide,True)

def thanks(prs,s):
    slide=prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide,BG_DARK); ribbon(slide,s.get('hero_style','blue'))
    add_text(slide,s.get('title','Thank you'),.8,2.05,7.8,.8,44,WHITE,True,'Arial Black')
    add_text(slide,s.get('subtitle',''),.88,3.05,6.6,.5,17,MUTED)
    add_text(slide,s.get('contact',''),.9,5.7,6.5,.35,13,WHITE,True)
    add_footer(slide,True)

LAYOUTS={
 'title':title_slide,'section_divider':section_divider,'section_word':section_word,'agenda':agenda,'feature_pill':feature_pill,'three_cards':three_cards,'pipeline':pipeline,'kpi':kpi,'quote':quote,'two_col':two_col,'bullets':bullets,'thanks':thanks
}

SAMPLE={
 "title":"Azure AI Foundry: Build Agents That Ship",
 "slides":[
  {"layout":"title","title":"Azure AI Foundry: Build Agents That Ship","subtitle":"MCAPS-style customer briefing deck","presenters":"Roey Zalta · Microsoft Israel","hero_style":"blue"},
  {"layout":"agenda","subtitle":"A practical briefing for startup engineering leaders","items":["What changed in the Azure AI platform","Reference architecture for production agents","Evaluation, safety, and observability","Migration path and next steps"]},
  {"layout":"feature_pill","badge":"Public preview","title":"Agentic apps need a control plane","subtitle":"Move from impressive demos to governed, observable production systems.","bullets":["Centralized model and tool routing","Built-in eval loops before release","Traceability from user request to action"]},
  {"layout":"three_cards","title":"Three pillars for production","subtitle":"The MCAPS conversation: value, architecture, and trust.","cards":[{"title":"Build faster","body":"Start from reusable agent patterns instead of hand-rolled orchestration."},{"title":"Operate safely","body":"Use identity, policy, content filters, and audit trails by default."},{"title":"Measure impact","body":"Tie quality, latency, and cost to business KPIs."}]},
  {"layout":"pipeline","title":"Reference implementation path","subtitle":"A clean path from first use case to production rollout.","steps":[{"icon":"1","title":"Discover","body":"Map workflows and user journeys"},{"icon":"2","title":"Ground","body":"Connect data, APIs, and actions"},{"icon":"3","title":"Evaluate","body":"Run golden sets and red-team tests"},{"icon":"4","title":"Deploy","body":"Ship with monitoring and rollback"},{"icon":"5","title":"Improve","body":"Close the loop with telemetry"}]},
  {"layout":"kpi","title":"Executive success metrics","stats":[{"value":"<50ms","label":"RAG fast path","note":"Cached or extractive answer path"},{"value":"99.9%","label":"Observed uptime","note":"Production SLO target"},{"value":"30%","label":"Cost reduction","note":"Via routing and caching"},{"value":"2 weeks","label":"Pilot window","note":"From discovery to demo"}]},
  {"layout":"thanks","title":"Thank you","subtitle":"Let’s build the first production-grade agent workflow.","contact":"roey.zalta@microsoft.com","hero_style":"purple"}
 ]
}

def build(outline, out):
    prs=Presentation(); prs.slide_width=I(W); prs.slide_height=I(H)
    for slide_spec in outline.get('slides',[]):
        layout=slide_spec.get('layout','bullets')
        fn=LAYOUTS.get(layout, bullets)
        fn(prs, slide_spec)
    prs.save(out)

if __name__=='__main__':
    if len(sys.argv)==1:
        print(json.dumps(SAMPLE, indent=2, ensure_ascii=False)); sys.exit(0)
    if len(sys.argv)!=3:
        print(__doc__); sys.exit(2)
    outline=json.loads(Path(sys.argv[1]).read_text())
    build(outline, sys.argv[2])
    print(f"wrote {sys.argv[2]}")
